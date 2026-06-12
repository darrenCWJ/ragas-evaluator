"""Text extraction from uploaded document files, one parser per format.

Each parser takes raw bytes and returns extracted text. ``parse_document``
dispatches on file extension; unsupported or unparsable files raise
``DocumentParseError`` with a user-facing message.
"""

from __future__ import annotations

import csv
import io
import json


class DocumentParseError(Exception):
    """Raised when a document cannot be parsed; message is user-facing."""


def _decode(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")


def _parse_text(data: bytes) -> str:
    return _decode(data)


def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _table_rows_to_markdown(rows: list[list[str]]) -> str:
    """Render rows as a GitHub-style markdown table (first row = header)."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    padded = [r + [""] * (width - len(r)) for r in rows]
    clean = [[" ".join(str(c).split()) for c in r] for r in padded]
    header, *body = clean
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(r) + " |" for r in body)
    return "\n".join(lines)


def _parse_docx(data: bytes, extract_tables: bool = True) -> str:
    """Paragraphs AND tables, in document order (tables become markdown)."""
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            parts.append(Paragraph(child, doc).text)
        elif isinstance(child, CT_Tbl) and extract_tables:
            rows = [
                [cell.text for cell in row.cells]
                for row in Table(child, doc).rows
            ]
            md = _table_rows_to_markdown(rows)
            if md:
                parts.append("\n" + md + "\n")
    return "\n".join(parts)


def _parse_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(data), "html.parser")
    for tag in soup(["script", "style", "nav", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_pptx(data: bytes, extract_tables: bool = True) -> str:
    """Slide text frames AND slide tables (tables become markdown)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            elif getattr(shape, "has_table", False) and extract_tables:
                rows = [
                    [cell.text for cell in row.cells]
                    for row in shape.table.rows
                ]
                md = _table_rows_to_markdown(rows)
                if md:
                    texts.append(md)
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


_XLSX_MAX_ROWS_PER_SHEET = 2000


def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        lines: list[str] = []
        for row in ws.iter_rows(max_row=_XLSX_MAX_ROWS_PER_SHEET, values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))
        if lines:
            parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(lines))
    wb.close()
    return "\n\n".join(parts)


def _parse_csv(data: bytes) -> str:
    text = _decode(data)
    lines = text.splitlines()
    delim = "\t" if lines and "\t" in lines[0] else ","
    reader = csv.reader(io.StringIO(text), delimiter=delim)
    return "\n".join(" | ".join(row) for row in reader if any(c.strip() for c in row))


def _parse_json(data: bytes) -> str:
    parsed = json.loads(_decode(data))
    return json.dumps(parsed, indent=2, ensure_ascii=False)


def _parse_yaml(data: bytes) -> str:
    # YAML is already human-readable text — validate it parses, store as-is.
    import yaml

    text = _decode(data)
    yaml.safe_load(text)
    return text


_PARSERS = {
    ".txt": _parse_text,
    ".md": _parse_text,
    ".markdown": _parse_text,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".csv": _parse_csv,
    ".tsv": _parse_csv,
    ".json": _parse_json,
    ".yaml": _parse_yaml,
    ".yml": _parse_yaml,
}

SUPPORTED_TEXT_EXTENSIONS = frozenset(_PARSERS)


# --- Embedded images (optional vision description) --------------------------

_IMAGE_MIME_BY_EXT = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "webp": "image/webp",
    "gif": "image/gif",
}
MAX_EMBEDDED_IMAGES = 10
_MIN_EMBEDDED_IMAGE_BYTES = 4 * 1024  # skip icons/bullets/decorations


def _image_mime(name: str) -> str | None:
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    return _IMAGE_MIME_BY_EXT.get(ext)


def extract_embedded_images(ext: str, data: bytes) -> list[dict]:
    """Embedded images from PPTX/DOCX/PDF for optional vision description.

    Returns up to MAX_EMBEDDED_IMAGES entries of
    ``{"label": str, "mime": str, "data": bytes}``; unsupported formats and
    extraction failures return an empty list (callers treat images as
    best-effort enrichment, never a parse blocker).
    """
    images: list[dict] = []
    try:
        if ext == ".pptx":
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(io.BytesIO(data))
            for i, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                    blob = shape.image.blob
                    mime = shape.image.content_type
                    if len(blob) < _MIN_EMBEDDED_IMAGE_BYTES or not mime.startswith("image/"):
                        continue
                    images.append({"label": f"Slide {i}", "mime": mime, "data": blob})
                    if len(images) >= MAX_EMBEDDED_IMAGES:
                        return images
        elif ext == ".docx":
            from docx import Document

            doc = Document(io.BytesIO(data))
            for rel in doc.part.rels.values():
                if "image" not in rel.reltype or rel.is_external:
                    continue
                blob = rel.target_part.blob
                mime = getattr(rel.target_part, "content_type", "") or (
                    _image_mime(rel.target_part.partname or "") or ""
                )
                if len(blob) < _MIN_EMBEDDED_IMAGE_BYTES or not mime.startswith("image/"):
                    continue
                images.append({"label": "Document image", "mime": mime, "data": blob})
                if len(images) >= MAX_EMBEDDED_IMAGES:
                    return images
        elif ext == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            for page_no, page in enumerate(reader.pages, start=1):
                for img in page.images:
                    mime = _image_mime(img.name or "")
                    if mime is None or len(img.data) < _MIN_EMBEDDED_IMAGE_BYTES:
                        continue
                    images.append({"label": f"Page {page_no}", "mime": mime, "data": img.data})
                    if len(images) >= MAX_EMBEDDED_IMAGES:
                        return images
    except Exception:  # noqa: BLE001 — enrichment only, never block the upload
        import logging

        logging.getLogger(__name__).warning(
            "Embedded-image extraction failed for %s", ext, exc_info=True
        )
    return images


def parse_document(filename: str, ext: str, data: bytes, *, extract_tables: bool = True) -> str:
    """Extract text from an uploaded document. Raises DocumentParseError.

    ``extract_tables`` (DOCX/PPTX) renders embedded tables as markdown rows
    inside the text; XLSX/CSV are inherently tabular and always extracted.
    """
    parser = _PARSERS.get(ext)
    if parser is None:
        raise DocumentParseError(f"Unsupported file type: {ext}")
    try:
        if ext == ".docx":
            text = _parse_docx(data, extract_tables)
        elif ext == ".pptx":
            text = _parse_pptx(data, extract_tables)
        else:
            text = parser(data)
    except DocumentParseError:
        raise
    except Exception as e:
        raise DocumentParseError(f"Could not read {ext.lstrip('.').upper()} file '{filename}': {e}") from e
    if not text.strip():
        raise DocumentParseError(
            f"No text could be extracted from '{filename}'. "
            "If this is a scanned/image PDF, upload page images instead (vision extraction)."
        )
    return text
