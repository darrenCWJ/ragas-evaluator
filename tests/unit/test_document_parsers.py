"""Document parser registry + vision model selection."""

import io
import json

import pytest
from fastapi import HTTPException

from pipeline.document_parsers import (
    SUPPORTED_TEXT_EXTENSIONS,
    DocumentParseError,
    parse_document,
)


class TestTextFormats:
    def test_plain_text(self):
        assert parse_document("a.txt", ".txt", b"hello world") == "hello world"

    def test_markdown_stored_verbatim(self):
        md = b"# Title\n\nSome **bold** text."
        assert parse_document("a.md", ".md", md) == md.decode()

    def test_json_pretty_printed(self):
        out = parse_document("a.json", ".json", b'{"a":1,"b":[2,3]}')
        assert json.loads(out) == {"a": 1, "b": [2, 3]}
        assert "\n" in out

    def test_yaml_validated_and_kept(self):
        out = parse_document("a.yaml", ".yaml", b"key: value\nitems:\n  - one\n")
        assert "key: value" in out

    def test_invalid_json_raises_user_facing_error(self):
        with pytest.raises(DocumentParseError, match="JSON"):
            parse_document("bad.json", ".json", b"{not json")

    def test_csv_to_pipe_table(self):
        out = parse_document("a.csv", ".csv", b"name,age\nalice,30\nbob,25\n")
        assert out.splitlines() == ["name | age", "alice | 30", "bob | 25"]

    def test_tsv_detected(self):
        out = parse_document("a.tsv", ".tsv", b"name\tage\nalice\t30\n")
        assert out.splitlines()[0] == "name | age"

    def test_html_strips_tags_and_scripts(self):
        html = b"<html><head><script>evil()</script></head><body><h1>Title</h1><p>Body text</p></body></html>"
        out = parse_document("a.html", ".html", html)
        assert "Title" in out
        assert "Body text" in out
        assert "evil" not in out


class TestOfficeFormats:
    def test_pptx_slide_markers(self):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        buf = io.BytesIO()
        prs.save(buf)
        out = parse_document("deck.pptx", ".pptx", buf.getvalue())
        assert "[Slide 1]" in out
        assert "Slide Title" in out

    def test_xlsx_sheet_rows(self):
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["name", "score"])
        ws.append(["alice", 0.9])
        buf = io.BytesIO()
        wb.save(buf)
        out = parse_document("book.xlsx", ".xlsx", buf.getvalue())
        assert "[Sheet: Data]" in out
        assert "alice | 0.9" in out


def _make_docx_with_table() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Quarterly results below.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Quarter"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "Q1"
    table.cell(1, 1).text = "1.2M"
    doc.add_paragraph("End of report.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_with_table() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Metrics"
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "Metric"
    shape.table.cell(0, 1).text = "Value"
    shape.table.cell(1, 0).text = "Recall"
    shape.table.cell(1, 1).text = "0.93"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestTableExtraction:
    def test_docx_tables_extracted_in_order(self):
        out = parse_document("report.docx", ".docx", _make_docx_with_table())
        assert "Quarterly results below." in out
        assert "| Quarter | Revenue |" in out
        assert "| Q1 | 1.2M |" in out
        # Document order preserved: intro before table before outro
        assert out.index("Quarterly") < out.index("| Q1 |") < out.index("End of report.")

    def test_docx_tables_can_be_disabled(self):
        out = parse_document(
            "report.docx", ".docx", _make_docx_with_table(), extract_tables=False
        )
        assert "Quarterly results below." in out
        assert "| Q1 |" not in out

    def test_pptx_slide_tables_extracted(self):
        out = parse_document("deck.pptx", ".pptx", _make_pptx_with_table())
        assert "[Slide 1]" in out
        assert "| Metric | Value |" in out
        assert "| Recall | 0.93 |" in out

    def test_pptx_tables_can_be_disabled(self):
        out = parse_document(
            "deck.pptx", ".pptx", _make_pptx_with_table(), extract_tables=False
        )
        assert "Metrics" in out  # title text still there
        assert "| Recall |" not in out


class TestEmbeddedImages:
    def test_unsupported_format_returns_empty(self):
        from pipeline.document_parsers import extract_embedded_images

        assert extract_embedded_images(".txt", b"hello") == []

    def test_docx_without_images_returns_empty(self):
        from pipeline.document_parsers import extract_embedded_images

        assert extract_embedded_images(".docx", _make_docx_with_table()) == []

    def test_corrupt_input_never_raises(self):
        from pipeline.document_parsers import extract_embedded_images

        assert extract_embedded_images(".pdf", b"not a real pdf") == []


class TestErrors:
    def test_unsupported_extension(self):
        with pytest.raises(DocumentParseError, match="Unsupported"):
            parse_document("a.exe", ".exe", b"binary")

    def test_empty_extraction_raises(self):
        with pytest.raises(DocumentParseError, match="No text"):
            parse_document("empty.txt", ".txt", b"   \n  ")

    def test_supported_extensions_match_config(self):
        from config import TEXT_FILE_TYPES

        assert set(SUPPORTED_TEXT_EXTENSIONS) == TEXT_FILE_TYPES


class TestVisionModelSelection:
    def test_explicit_vision_model_wins(self, monkeypatch):
        monkeypatch.setattr("config.VISION_MODEL", "gemini-2.0-flash")
        from pipeline.vision import pick_vision_model

        assert pick_vision_model() == "gemini-2.0-flash"

    def test_no_keys_raises_422(self, monkeypatch):
        monkeypatch.setattr("config.VISION_MODEL", "")
        for var in ("OPENAI_API_KEY", "ANTHROPIC_API_KEY", "GOOGLE_API_KEY", "OPENAI_BASE_URL"):
            monkeypatch.delenv(var, raising=False)
        from pipeline.vision import pick_vision_model

        with pytest.raises(HTTPException) as exc:
            pick_vision_model()
        assert exc.value.status_code == 422

    def test_openai_key_picks_gpt(self, monkeypatch):
        monkeypatch.setattr("config.VISION_MODEL", "")
        monkeypatch.delenv("OPENAI_BASE_URL", raising=False)
        monkeypatch.setenv("OPENAI_API_KEY", "sk-test")
        from pipeline.vision import pick_vision_model

        assert pick_vision_model() == "gpt-4o-mini"
